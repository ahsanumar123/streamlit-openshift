import streamlit as st
from pptx import Presentation
from pdf2image import convert_from_bytes
import plotly.graph_objects as go

st.title("Slide Viewer with Notes & Slide Image")

# Uploaders for PowerPoint (.pptx) and PDF files
uploaded_pptx = st.file_uploader("Upload a PowerPoint file (.pptx)", type=["pptx"])
uploaded_pdf = st.file_uploader("Upload the corresponding PDF file (.pdf)", type=["pdf"])

if uploaded_pptx and uploaded_pdf:
    # Extract slide notes from the PPTX file
    prs = Presentation(uploaded_pptx)
    slides = {}
    for i, slide in enumerate(prs.slides, start=1):
        slide_label = f"Slide {i}"
        if slide.has_notes_slide:
            # Access the notes text from the notes_text_frame property
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
        else:
            notes_text = "No notes available"
        slides[slide_label] = notes_text

    # Convert PDF pages to images (each page is assumed to represent a slide)
    pdf_bytes = uploaded_pdf.read()
    slide_images = convert_from_bytes(pdf_bytes)

    # Let user select a slide from a dropdown
    selected_slide = st.selectbox("Select a Slide", list(slides.keys()))
    slide_index = int(selected_slide.split()[1]) - 1

    # Retrieve the corresponding slide notes
    notes = slides[selected_slide]
    
    # Compute KPIs: word and character counts for the selected slide notes
    word_count = len(notes.split())
    char_count = len(notes)
    
    # Display KPIs as large numbers
    st.markdown("## Slide KPIs")
    col1, col2 = st.columns(2)
    col1.metric("Word Count", word_count)
    col2.metric("Character Count", char_count)
    
    # Create a simple Plotly bar chart to visualize the metrics
    fig = go.Figure(data=[go.Bar(x=["Words", "Characters"], y=[word_count, char_count])])
    fig.update_layout(title_text="Notes Metrics", xaxis_title="Metric", yaxis_title="Count")
    st.plotly_chart(fig)
    
    # Display the slide notes text
    st.markdown("## Slide Notes")
    st.write(notes)
    
    # Display the corresponding slide image from the PDF
    st.markdown("## Slide Image")
    if slide_index < len(slide_images):
        st.image(slide_images[slide_index], caption=selected_slide)
    else:
        st.warning("Corresponding slide image not found in the uploaded PDF.")

